from llama_index.llms.openai import OpenAI
from pptx import Presentation
import re
import json


class PowerpointGenerator:
    def init(self) -> None:
        self._init_message()

    def _init_message(self) -> None:
        print("Initializing Powerpoint Generator...")

    def _extract_json(self, slides_response: str) -> dict:
        print("Extracting JSON...")
        # Searching for JSON in the response
        json_match = re.search(r'```json\s*(.*?)\s*```', slides_response, re.DOTALL)

        # Verification and parsing of JSON
        if json_match:
            json_content = json_match.group(1)
            try:
                slides = json.loads(json_content)
            except json.JSONDecodeError as e:
                raise ValueError(f"JSON parsing error: {str(e)}")
        else:
            raise ValueError("The JSON was not found in the model response.")
        
        return slides

    def _populate_slides(self, slide: dict, slide_data: dict) -> None:
        # Get the list of placeholders
        placeholders = slide_data.get("placeholders", [])
        
        # Populate text placeholders
        for placeholder in placeholders:
            for shape in slide.shapes:
                if shape.has_text_frame and placeholder.get("name") in shape.name:
                    print(f"Populating text shape: {shape.name}")
                    if placeholder.get("text") and shape.has_text_frame:
                        shape.text = placeholder["text"]

        # Populate image placeholders
        # for placeholder in placeholders:
        #     if placeholder.get("image_description"):
        #         image_data = generate_image(placeholder["image_description"])
        #         img_path = "temp_image.png"
        #         with open(img_path, "wb") as img_file:
        #             img_file.write(image_data)
        #     for shape in slide.shapes:
        #         if shape.is_placeholder and shape.placeholder_format.type == 18:
        #             slide.shapes._spTree.remove(shape._element)
        #             slide.shapes.add_picture(img_path, shape.left, shape.top, shape.width, shape.height)
        #             break

        # Populate the tables
        for placeholder in placeholders:
            table_data = placeholder.get("table_structure")
            if table_data:
                print("Populating table...")
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for i, row_data in enumerate(table_data):
                            for j, cell_text in enumerate(row_data):
                                table.cell(i, j).text = cell_text
                    break
    
    def _populate_all_slides(self, prs: Presentation, slides: dict) -> None:
        print("Populating slides...")
        for i, slide_data in enumerate(slides):
            if i < len(prs.slides):
                print(f"Populating slide {i+1}...")
                self._populate_slides(prs.slides[i], slide_data)
    
    def generate(self, prs: Presentation, content_prompt: str, llm: OpenAI, output_path: str) -> None:
        print("Generating slides...")
        try:
            slides_response = llm.complete(content_prompt, True).text
            slides = self._extract_json(slides_response)
            self._populate_all_slides(prs, slides)
            prs.save(output_path)
            print("Slides generated successfully. File Saved.")
        except Exception as e:
            print(f"Error: {str(e)}")
        