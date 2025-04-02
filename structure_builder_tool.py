from pptx import Presentation
import json

class StructureBuilderTool:
    def __init__(self) -> None:
        self.slide_structure = []
        self._init_message()
    
    def _init_message(self) -> None:
        print("Initializing Structure Builder...")

    def build_structure(self, prs: Presentation) -> str:
        print("Building structure...")
        for slide in prs.slides:
            structure = {"placeholders": []}
            
            # Iterate through each shape in the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Create a dictionary to store placeholder details
                    text_shape = {
                        "type": "text_box" if not shape.is_placeholder else "placeholder",
                        "idx": shape.placeholder_format.idx if shape.is_placeholder else None,
                        "has_text_frame": True,
                        "name": shape.name,
                        "text": shape.text_frame.text if shape.text_frame else ""
                    }
                    if shape.is_placeholder:
                        text_shape["placeholder_type"] = shape.placeholder_format.type
                        # Check if the placeholder contains a table
                        if shape.has_table:
                            text_shape["has_table"] = True
                            # Extract the table structure as a list of rows with cell texts
                            text_shape["table_structure"] = [
                                [cell.text for cell in row.cells] for row in shape.table.rows
                            ]
                        else:
                            text_shape["has_table"] = False
    
                        # Check if the placeholder is an image placeholder
                        if shape.placeholder_format.type == 18:
                            text_shape["has_image"] = True
                            text_shape["image_description"] = ""
                        else:
                            text_shape["has_image"] = False

                    # Append the placeholder information to the structure
                    structure["placeholders"].append(text_shape)

            # Append the slide structure to the list of all slide structures
            self.slide_structure.append(structure)
    
        # Convert the structure list to a JSON string for better readability
        slide_structures_str = json.dumps(self.slide_structure, indent=4)
        return slide_structures_str